import httpx
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt

async def handle_ppt(topic: str) -> str:
    async with httpx.AsyncClient() as client:
        response = await client.post(
            url="https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {os.getenv('GROQ_API_KEY')}",
                "Content-Type": "application/json",
            },
            json={
                "model": "llama-3.3-70b-versatile",
                "max_tokens": 1024,
                "response_format": {"type": "json_object"},
                "messages": [
                    {
                        "role": "user",
                        "content": f"""Create a 5 slide presentation about: {topic}
Return ONLY this JSON format:
{{
    "title": "presentation title",
    "slides": [
        {{
            "title": "slide title",
            "points": ["point 1", "point 2", "point 3"]
        }}
    ]
}}""",
                    }
                ],
            },
        )

    content = response.json()["choices"][0]["message"]["content"]
    slides_data = json.loads(content)

    # Create PPT
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = slides_data["title"]

    # Content slides
    for slide in slides_data["slides"]:
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        content_slide.shapes.title.text = slide["title"]

        # Add bullet points
        tf = content_slide.placeholders[1].text_frame
        for i, point in enumerate(slide["points"]):
            if i == 0:
                tf.text = point
            else:
                tf.add_paragraph().text = point


# Save file
    filename = f"{slides_data['title'].replace(' ', '_')}.pptx"
    filepath = f"outputs/{filename}"
    os.makedirs("outputs", exist_ok=True)
    prs.save(filepath)

    base_url = os.getenv("RAILWAY_PUBLIC_DOMAIN", "localhost:8000")
    download_url = f"https://{base_url}/download/{filename}"
    return f"PPT created successfully! [Download your presentation]({download_url})"